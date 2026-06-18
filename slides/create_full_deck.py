#!/usr/bin/env python3
"""
Create full 2-slide deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLACK = RGBColor(25, 25, 25)
ACCENT = RGBColor(0, 82, 147)
GRAY = RGBColor(120, 120, 120)
WHITE = RGBColor(255, 255, 255)

# ==================== SLIDE 1: CASCADE EFFECTS ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
bg.line.fill.background()

title_box = slide1.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(10), Inches(0.5))
tf = title_box.text_frame
tf.paragraphs[0].text = "Hormuz Reopening: Industry Impact"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

sub_box = slide1.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(10), Inches(0.3))
tf = sub_box.text_frame
tf.paragraphs[0].text = "First, Second, and Third Order Effects"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"

col_width = Inches(3.5)
col_x = [Inches(0.7), Inches(4.5), Inches(8.3)]
header_y = Inches(1.7)
content_y = Inches(2.3)

headers = [("1ST ORDER", "Immediate"), ("2ND ORDER", "1-3 Months"), ("3RD ORDER", "3-12 Months")]

for i, (title, timing) in enumerate(headers):
    label_box = slide1.shapes.add_textbox(col_x[i], header_y, col_width, Inches(0.3))
    tf = label_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.name = "Arial"

    time_box = slide1.shapes.add_textbox(col_x[i] + Inches(1.5), header_y, Inches(1.5), Inches(0.3))
    tf = time_box.text_frame
    tf.paragraphs[0].text = timing
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = "Arial"

    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x[i], header_y + Inches(0.35), Inches(3.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

first_order = [
    ("Oil & Gas", "War premium unwinds; Brent down 5%", "▼ 4%"),
    ("Shipping", "Routes reopening; insurance still elevated", "▲ 3%"),
    ("Airlines", "Jet fuel costs falling (31% of opex)", "▲ 4%"),
    ("Refiners", "Input costs falling; margins expand", "▲ 3%"),
]
second_order = [
    ("Consumer", "~$40/mo back in pockets from gas", "▲ 2%"),
    ("Materials", "Petrochemical feedstock costs down", "▲ 3%"),
    ("Retail", "Lower logistics; more spending", "▲ 2%"),
    ("Autos", "Parts flow resuming; costs normalize", "▲ 2%"),
]
third_order = [
    ("Semiconductors", "Helium recovery takes years", "— Flat"),
    ("Agriculture", "Fertilizer easing; contracts lag", "▲ 2%"),
    ("Inflation", "Energy deflation feeds to core", "▼ 0.5%"),
    ("Fed Policy", "Lower CPI could enable cuts", "— TBD"),
]

all_content = [first_order, second_order, third_order]

for col_idx, items in enumerate(all_content):
    y = content_y
    for name, detail, metric in items:
        name_box = slide1.shapes.add_textbox(col_x[col_idx], y, Inches(2.2), Inches(0.25))
        tf = name_box.text_frame
        tf.paragraphs[0].text = name
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLACK
        tf.paragraphs[0].font.name = "Arial"

        metric_box = slide1.shapes.add_textbox(col_x[col_idx] + Inches(2.2), y, Inches(1), Inches(0.25))
        tf = metric_box.text_frame
        tf.paragraphs[0].text = metric
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        if "▲" in metric:
            tf.paragraphs[0].font.color.rgb = ACCENT
        elif "▼" in metric:
            tf.paragraphs[0].font.color.rgb = BLACK
        else:
            tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].font.name = "Arial"
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        detail_box = slide1.shapes.add_textbox(col_x[col_idx], y + Inches(0.25), Inches(3.2), Inches(0.4))
        tf = detail_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = detail
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].font.name = "Arial"

        y += Inches(0.85)

# Flow arrows
arrow_y = Inches(3.8)
for i in range(2):
    arrow = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     col_x[i] + col_width + Inches(0.15), arrow_y,
                                     Inches(0.6), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()

# Footnote explaining percentages
footnote = slide1.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(12), Inches(0.3))
tf = footnote.text_frame
tf.paragraphs[0].text = "% = Stock price moves on June 15, 2026 (XLE, DAL, UAL data). Sources: Reuters, CNBC, Yahoo Finance"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"

# Immediate impact callout box
callout_bg = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.7))
callout_bg.fill.solid()
callout_bg.fill.fore_color.rgb = RGBColor(240, 247, 255)
callout_bg.line.color.rgb = ACCENT
callout_bg.line.width = Pt(1)

key_label = slide1.shapes.add_textbox(Inches(0.9), Inches(6.3), Inches(2.5), Inches(0.25))
tf = key_label.text_frame
tf.paragraphs[0].text = "IMMEDIATE IMPACT:"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = ACCENT
tf.paragraphs[0].font.name = "Arial"

key_text = slide1.shapes.add_textbox(Inches(3.2), Inches(6.3), Inches(9.2), Inches(0.5))
tf = key_text.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Airlines (+4%), Shipping (+3%), Refiners (+3%) benefit now. Oil & Gas (-4%) faces headwinds as war premium unwinds."
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

footer = slide1.shapes.add_textbox(Inches(0.7), Inches(7.1), Inches(5), Inches(0.25))
tf = footer.text_frame
tf.paragraphs[0].text = "June 2026"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"


# ==================== SLIDE 2: PLAIN ENGLISH SUMMARY ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg2.fill.solid()
bg2.fill.fore_color.rgb = WHITE
bg2.line.fill.background()

title_box2 = slide2.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
tf = title_box2.text_frame
tf.paragraphs[0].text = "What This Means for Business"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

sub_box2 = slide2.shapes.add_textbox(Inches(0.7), Inches(1.05), Inches(12), Inches(0.4))
tf = sub_box2.text_frame
tf.paragraphs[0].text = "The peace deal changes who wins and who loses"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"

content = [
    ("Cheaper to fly and ship",
     "Oil prices are down 5% since the deal was announced. Airlines and shipping companies will pay less for fuel, which means lower costs and higher profits."),
    ("More money in consumers' pockets",
     "Gas prices falling from $4.56 to $4.12 per gallon. Average household could save around $40/month as prices continue dropping."),
    ("Supply chains may start to ease",
     "Hormuz expected to reopen within 30 days, but backlogs and elevated shipping costs will take months to clear."),
    ("Tech still has a problem",
     "Chip makers need helium, and Qatar's production was hit by strikes. Recovery takes years, not months."),
    ("Oil companies lose their advantage",
     "Energy stocks fell 3.75% on June 15. High oil prices helped US drillers. Now that prices are falling, the windfall is ending."),
]

y = Inches(1.7)
for headline, detail in content:
    bullet = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.08), Inches(0.12), Inches(0.12))
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = ACCENT
    bullet.line.fill.background()

    head_box = slide2.shapes.add_textbox(Inches(1.0), y, Inches(11), Inches(0.35))
    tf = head_box.text_frame
    tf.paragraphs[0].text = headline
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK
    tf.paragraphs[0].font.name = "Arial"

    detail_box = slide2.shapes.add_textbox(Inches(1.0), y + Inches(0.35), Inches(11), Inches(0.6))
    tf = detail_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = detail
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = "Arial"

    y += Inches(0.95)

line2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(6.7), Inches(12), Inches(0.02))
line2.fill.solid()
line2.fill.fore_color.rgb = ACCENT
line2.line.fill.background()

takeaway = slide2.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(12), Inches(0.5))
tf = takeaway.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Bottom line: Travel, retail, and shipping benefit now. Manufacturing recovers over months. Tech stays constrained."
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

# Save
output_path = "/Users/jamesvandehei/Documents/iran-metrics-dashboard/slides/Hormuz_Cascade_Effects.pptx"
prs.save(output_path)
print(f"Saved {len(prs.slides)} slides to: {output_path}")
