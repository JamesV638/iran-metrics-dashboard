#!/usr/bin/env python3
"""
Create clean, minimal PowerPoint slide - 2 colors only
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# === COLORS (2 only) ===
BLACK = RGBColor(25, 25, 25)
ACCENT = RGBColor(0, 82, 147)  # Deep blue
GRAY = RGBColor(120, 120, 120)
LIGHT_GRAY = RGBColor(245, 245, 245)
WHITE = RGBColor(255, 255, 255)

# === BACKGROUND ===
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
bg.line.fill.background()

# === TITLE ===
title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(10), Inches(0.5))
tf = title_box.text_frame
tf.paragraphs[0].text = "Hormuz Reopening: Industry Impact"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

# === SUBTITLE ===
sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(10), Inches(0.3))
tf = sub_box.text_frame
tf.paragraphs[0].text = "First, Second, and Third Order Effects"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"

# === COLUMN SETUP ===
col_width = Inches(3.5)
col_x = [Inches(0.7), Inches(4.5), Inches(8.3)]
header_y = Inches(1.7)
content_y = Inches(2.3)

# === HEADERS ===
headers = [
    ("1ST ORDER", "Immediate"),
    ("2ND ORDER", "1-3 Months"),
    ("3RD ORDER", "3-12 Months")
]

for i, (title, timing) in enumerate(headers):
    # Order label
    label_box = slide.shapes.add_textbox(col_x[i], header_y, col_width, Inches(0.3))
    tf = label_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.name = "Arial"

    # Timing
    time_box = slide.shapes.add_textbox(col_x[i] + Inches(1.5), header_y, Inches(1.5), Inches(0.3))
    tf = time_box.text_frame
    tf.paragraphs[0].text = timing
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = "Arial"

    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x[i], header_y + Inches(0.35), Inches(3.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

# === CONTENT ===
first_order = [
    ("Oil & Gas", "War premium unwinds; Brent -23%", "▼ 5%"),
    ("Shipping", "Insurance drops 60-70%; routes reopen", "▲ 8%"),
    ("Airlines", "Jet fuel costs plummet (30% of opex)", "▲ 10%"),
    ("Refiners", "Input costs falling; margins expand", "▲ 6%"),
]

second_order = [
    ("Consumer", "$85/mo back in pockets from gas", "▲ 3%"),
    ("Materials", "Petrochemical feedstock costs down", "▲ 5%"),
    ("Retail", "Lower logistics; more spending", "▲ 5%"),
    ("Autos", "Parts flow resuming; costs normalize", "▲ 4%"),
]

third_order = [
    ("Semiconductors", "Helium recovery takes years", "— Flat"),
    ("Agriculture", "Fertilizer easing; contracts lag", "▲ 2%"),
    ("Inflation", "Energy deflation feeds to core", "▼ 0.5%"),
    ("Fed Policy", "Lower CPI could enable cuts", "— TBD"),
]

all_content = [first_order, second_order, third_order]

for col_idx, items in enumerate(all_content):
    y = content_y
    for name, detail, metric in items:
        # Industry name
        name_box = slide.shapes.add_textbox(col_x[col_idx], y, Inches(2.2), Inches(0.25))
        tf = name_box.text_frame
        tf.paragraphs[0].text = name
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLACK
        tf.paragraphs[0].font.name = "Arial"

        # Metric
        metric_box = slide.shapes.add_textbox(col_x[col_idx] + Inches(2.2), y, Inches(1), Inches(0.25))
        tf = metric_box.text_frame
        tf.paragraphs[0].text = metric
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        if "▲" in metric:
            tf.paragraphs[0].font.color.rgb = ACCENT
        elif "▼" in metric:
            tf.paragraphs[0].font.color.rgb = BLACK
        else:
            tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].font.name = "Arial"
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Detail
        detail_box = slide.shapes.add_textbox(col_x[col_idx], y + Inches(0.25), Inches(3.2), Inches(0.4))
        tf = detail_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = detail
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].font.name = "Arial"

        y += Inches(0.85)

# === FLOW ARROWS ===
arrow_y = Inches(3.8)
arrow_width = Inches(0.6)

# Arrow 1
arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 col_x[0] + col_width + Inches(0.15), arrow_y,
                                 arrow_width, Inches(0.25))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = ACCENT
arrow1.line.fill.background()

# Arrow 2
arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 col_x[1] + col_width + Inches(0.15), arrow_y,
                                 arrow_width, Inches(0.25))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = ACCENT
arrow2.line.fill.background()

# === KEY TAKEAWAY ===
key_y = Inches(6.4)

# Label
key_label = slide.shapes.add_textbox(Inches(0.7), key_y, Inches(1), Inches(0.25))
tf = key_label.text_frame
tf.paragraphs[0].text = "KEY:"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = ACCENT
tf.paragraphs[0].font.name = "Arial"

# Text
key_text = slide.shapes.add_textbox(Inches(1.3), key_y, Inches(11), Inches(0.5))
tf = key_text.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Airlines (+10%), Shipping (+8%), and Retail (+5%) see immediate benefit. Oil & Gas faces headwinds. Semiconductors remain constrained — helium shortage takes years to resolve."
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

# === FOOTER LINE ===
footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(7.0), Inches(12), Inches(0.01))
footer_line.fill.solid()
footer_line.fill.fore_color.rgb = LIGHT_GRAY
footer_line.line.fill.background()

# === FOOTER ===
footer = slide.shapes.add_textbox(Inches(0.7), Inches(7.1), Inches(5), Inches(0.25))
tf = footer.text_frame
tf.paragraphs[0].text = "June 2026"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"

# === SAVE ===
output_path = "/Users/jamesvandehei/Documents/iran-metrics-dashboard/slides/Hormuz_Cascade_Effects.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
