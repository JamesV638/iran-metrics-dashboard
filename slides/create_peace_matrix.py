#!/usr/bin/env python3
"""
Create 2x2 scenario matrix slide: "Will the peace hold?"
Matching FGS Global slide specs exactly
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation - 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLACK = RGBColor(30, 30, 30)
GRAY = RGBColor(120, 120, 120)
LIGHT_GRAY = RGBColor(200, 200, 200)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(180, 100, 80)  # Muted rust/salmon for headers
BORDER = RGBColor(180, 180, 180)

# Add slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# White background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
bg.line.fill.background()

# === HEADER ===
header = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.3))
tf = header.text_frame
tf.paragraphs[0].text = "JUNE 2026"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"
tf.paragraphs[0].font.bold = True

# === MAIN TITLE ===
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(8), Inches(0.7))
tf = title.text_frame
tf.paragraphs[0].text = "Will the peace hold?"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

# === LEFT SIDEBAR CONTEXT ===
context = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(2.0), Inches(2.5))
tf = context.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "How implementation progress and stakeholder alignment will shape whether the 14-point framework becomes a durable settlement or unravels. Four scenarios for businesses to consider."
p.font.size = Pt(11)
p.font.color.rgb = BLACK
p.font.name = "Arial"
p.line_spacing = 1.3

# === 2x2 MATRIX ===
matrix_left = Inches(2.8)
matrix_top = Inches(1.5)
cell_width = Inches(2.9)
cell_height = Inches(2.3)
gap = Inches(0.05)

# Top axis label
top_label = slide.shapes.add_textbox(matrix_left, Inches(1.2), cell_width * 2 + gap, Inches(0.25))
tf = top_label.text_frame
tf.paragraphs[0].text = "Unified Stakeholders"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Bottom axis label
bottom_label = slide.shapes.add_textbox(matrix_left, Inches(6.15), cell_width * 2 + gap, Inches(0.25))
tf = bottom_label.text_frame
tf.paragraphs[0].text = "Fragmented Stakeholders"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Left axis label (rotated text - using textbox positioned on left)
left_label = slide.shapes.add_textbox(Inches(2.35), Inches(2.8), Inches(0.4), Inches(2))
tf = left_label.text_frame
tf.paragraphs[0].text = "Stalled"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Right axis label
right_label = slide.shapes.add_textbox(Inches(8.55), Inches(2.8), Inches(0.6), Inches(2))
tf = right_label.text_frame
tf.paragraphs[0].text = "Advancing"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Matrix border
border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, matrix_left, matrix_top, cell_width * 2 + gap, cell_height * 2 + gap)
border.fill.background()
border.line.color.rgb = BORDER
border.line.width = Pt(1)

# Horizontal divider
h_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, matrix_left, matrix_top + cell_height, cell_width * 2 + gap, Pt(1))
h_line.fill.solid()
h_line.fill.fore_color.rgb = BORDER
h_line.line.fill.background()

# Vertical divider
v_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, matrix_left + cell_width, matrix_top, Pt(1), cell_height * 2 + gap)
v_line.fill.solid()
v_line.fill.fore_color.rgb = BORDER
v_line.line.fill.background()

# === QUADRANT CONTENT ===

def add_quadrant(x, y, title_text, desc, implications):
    # Title
    title_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), cell_width - Inches(0.2), Inches(0.3))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.name = "Arial"

    # Description
    desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.4), cell_width - Inches(0.2), Inches(0.9))
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = desc
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = BLACK
    tf.paragraphs[0].font.name = "Arial"
    tf.paragraphs[0].line_spacing = 1.2

    # Implications
    impl_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.3), cell_width - Inches(0.2), Inches(0.9))
    tf = impl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    # Bold "Implications" part
    run1 = p.add_run()
    run1.text = "Implications / key business questions: "
    run1.font.size = Pt(9)
    run1.font.bold = True
    run1.font.color.rgb = ACCENT
    run1.font.name = "Arial"

    run2 = p.add_run()
    run2.text = implications
    run2.font.size = Pt(9)
    run2.font.color.rgb = BLACK
    run2.font.name = "Arial"

# Top-Left: Negotiated Friction (Low implementation + Unified)
add_quadrant(
    matrix_left, matrix_top,
    "Negotiated Friction",
    "Leadership alignment holds, but 14-point implementation slows. Bureaucratic delays, disputed terms, or domestic politics create friction. Progress is real but erratic.",
    "How long should we maintain crisis-level inventory? Which interim shipping arrangements remain necessary? How do we plan when timelines keep shifting?"
)

# Top-Right: Durable Settlement (High implementation + Unified)
add_quadrant(
    matrix_left + cell_width + gap, matrix_top,
    "Durable Settlement",
    "Both sides deliver on the 14 points. Hardliner opposition contained. Hormuz fully operational within 30 days, sanctions relief proceeds, nuclear talks advance on schedule.",
    "When can we treat supply chains as normalized? How quickly should we unwind hedges and crisis premiums? What new opportunities emerge in post-conflict markets?"
)

# Bottom-Left: Deal Collapse (Low implementation + Fragmented)
add_quadrant(
    matrix_left, matrix_top + cell_height + gap,
    "Deal Collapse",
    "Framework unravels. Nuclear talks fail, disputed terms trigger walkouts, hardliners prevail on both sides. Return to blockade, conflict escalation, or worse.",
    "What is our exposure to renewed disruption? Are crisis contingency plans still current? How quickly can we pivot to alternative routes, suppliers, or markets?"
)

# Bottom-Right: Fragile Compliance (High implementation + Fragmented)
add_quadrant(
    matrix_left + cell_width + gap, matrix_top + cell_height + gap,
    "Fragile Compliance",
    "Both sides meet formal commitments, but spoilers remain active. IRGC hardliners, Israeli opposition, or disputed reconstruction terms could trigger incidents. Compliance without conviction.",
    "What early warning indicators should we monitor? How do we distinguish noise from signal? What contingency triggers should we pre-define for rapid response?"
)

# === RIGHT PANEL ===
right_panel_x = Inches(9.2)

# Right panel title
right_title = slide.shapes.add_textbox(right_panel_x, Inches(0.6), Inches(3.8), Inches(0.5))
tf = right_title.text_frame
tf.paragraphs[0].text = "What the Markets Are Pricing"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

# Subtitle
right_sub = slide.shapes.add_textbox(right_panel_x, Inches(1.0), Inches(3.8), Inches(0.4))
tf = right_sub.text_frame
tf.paragraphs[0].text = "Permanent peace deal probability"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"

# Placeholder for chart area (simple box representation)
chart_area = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_panel_x, Inches(1.5), Inches(3.6), Inches(3.8))
chart_area.fill.solid()
chart_area.fill.fore_color.rgb = RGBColor(250, 250, 250)
chart_area.line.color.rgb = LIGHT_GRAY
chart_area.line.width = Pt(0.5)

# Y-axis labels
for i, (val, y_pos) in enumerate([(100, 1.6), (75, 2.55), (50, 3.5), (25, 4.45), (0, 5.4)]):
    label = slide.shapes.add_textbox(right_panel_x - Inches(0.3), Inches(y_pos), Inches(0.3), Inches(0.2))
    tf = label.text_frame
    tf.paragraphs[0].text = str(val)
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = "Arial"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

# Chart note
chart_note = slide.shapes.add_textbox(right_panel_x, Inches(5.4), Inches(3.6), Inches(0.8))
tf = chart_note.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Polymarket odds for permanent US-Iran peace deal have risen from 15% pre-framework to 45%+ post-signing, reflecting cautious optimism about implementation."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"
tf.paragraphs[0].font.italic = True

# === FOOTER ===
# Footer line
footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.0), Inches(12.3), Pt(0.5))
footer_line.fill.solid()
footer_line.fill.fore_color.rgb = LIGHT_GRAY
footer_line.line.fill.background()

# Footer text left
footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.25))
tf = footer_left.text_frame
tf.paragraphs[0].text = "ALTERNATIVE FUTURES"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"
tf.paragraphs[0].font.bold = True

# Footer text right
footer_right = slide.shapes.add_textbox(Inches(10.5), Inches(7.1), Inches(2.3), Inches(0.25))
tf = footer_right.text_frame
tf.paragraphs[0].text = "Data: Polymarket, June 2026"
tf.paragraphs[0].font.size = Pt(8)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"
tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

# Page number
page_num = slide.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.3), Inches(0.25))
tf = page_num.text_frame
tf.paragraphs[0].text = "1"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"

# Save
output_path = "/Users/jamesvandehei/Documents/iran-metrics-dashboard/slides/Peace_Matrix_Scenarios.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
