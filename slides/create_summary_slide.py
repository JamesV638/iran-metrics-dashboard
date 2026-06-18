#!/usr/bin/env python3
"""
Create plain English summary slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Open existing presentation
prs = Presentation("/Users/jamesvandehei/Documents/iran-metrics-dashboard/slides/Hormuz_Cascade_Effects.pptx")

# Add new slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Colors
BLACK = RGBColor(25, 25, 25)
ACCENT = RGBColor(0, 82, 147)
GRAY = RGBColor(120, 120, 120)
WHITE = RGBColor(255, 255, 255)

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
bg.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
tf = title_box.text_frame
tf.paragraphs[0].text = "What This Means for Business"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.05), Inches(12), Inches(0.4))
tf = sub_box.text_frame
tf.paragraphs[0].text = "The peace deal changes who wins and who loses"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = GRAY
tf.paragraphs[0].font.name = "Arial"

# Main content - simple bullet points
content = [
    ("Cheaper to fly and ship",
     "Oil prices are dropping fast. Airlines and shipping companies will pay less for fuel, which means lower costs and higher profits."),

    ("More money in consumers' pockets",
     "Gas prices are falling. The average household could save $85/month — money that goes back into restaurants, retail, and travel."),

    ("Factories can get supplies again",
     "Ships can now pass through the Strait of Hormuz. Parts and materials that were stuck are moving again, helping manufacturers catch up."),

    ("Tech still has a problem",
     "Chip makers need helium, and that shortage won't be fixed for years. Peace deal or not, electronics and AI companies still face constraints."),

    ("Oil companies lose their advantage",
     "High oil prices helped US drillers. Now that prices are falling, their windfall is ending."),
]

y = Inches(1.7)
for headline, detail in content:
    # Bullet point
    bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.08), Inches(0.12), Inches(0.12))
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = ACCENT
    bullet.line.fill.background()

    # Headline
    head_box = slide.shapes.add_textbox(Inches(1.0), y, Inches(11), Inches(0.35))
    tf = head_box.text_frame
    tf.paragraphs[0].text = headline
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK
    tf.paragraphs[0].font.name = "Arial"

    # Detail
    detail_box = slide.shapes.add_textbox(Inches(1.0), y + Inches(0.35), Inches(11), Inches(0.6))
    tf = detail_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = detail
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = "Arial"

    y += Inches(0.95)

# Bottom line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(6.7), Inches(12), Inches(0.02))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Takeaway
takeaway = slide.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(12), Inches(0.5))
tf = takeaway.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Bottom line: Travel, retail, and shipping benefit now. Manufacturing recovers over months. Tech stays constrained."
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
tf.paragraphs[0].font.name = "Arial"

# Save
output_path = "/Users/jamesvandehei/Documents/iran-metrics-dashboard/slides/Hormuz_Cascade_Effects.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
